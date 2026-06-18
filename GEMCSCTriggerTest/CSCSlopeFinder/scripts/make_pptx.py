"""
Build a PowerPoint with 4 images per slide:
  - top-left:  endcap1 even
  - top-right: endcap2 even
  - bot-left:  endcap1 odd
  - bot-right: endcap2 odd

Run after the plotting script has produced PNGs in `plotdir`.
"""

import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

# ------------------------------------------------------------------
# Configuration
# ------------------------------------------------------------------
plotdir = "plots/alignment_both_aligned_data/"
output_pptx = os.path.join(plotdir, "alignment_plots.pptx")

charges = ["negative", "positive"]
endcaps = [1, 2]
parities = ["even", "odd"]

plot_families = [
    ("slope vs chamber",
     "bendingangevschamberendcap{endcap}_{charge}.png", False),
    ("|eta| vs slope",
     "eta{parity}chambersendcap{endcap}_{charge}.png", True),
    ("Muon pt",
     "ptplot{parity}endcap{endcap}_{charge}.png", True),
    ("pt*charge vs chamber",
     "chargexptendcap{endcap}_{charge}.png", False),
    ("pt*charge vs slope",
     "chargedpt_vs_chargedslope_{parity}_endcap{endcap}_{charge}.png", True),
]


# ------------------------------------------------------------------
# Helpers
# ------------------------------------------------------------------
def add_title_slide(prs, title_text, subtitle_text=""):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title_text
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle_text
    return slide


def add_picture_fit(slide, path, x, y, cell_w, cell_h):
    """Add a picture preserving aspect ratio, centered in the (x,y,cell_w,cell_h) box."""
    try:
        with Image.open(path) as im:
            iw, ih = im.size
    except Exception:
        # Fallback: just stick it in (shouldn't usually happen)
        slide.shapes.add_picture(path, x, y, width=cell_w, height=cell_h)
        return

    img_aspect = iw / ih
    cell_aspect = cell_w / cell_h

    if img_aspect > cell_aspect:
        # image is wider -> fit to width
        new_w = cell_w
        new_h = int(cell_w / img_aspect)
    else:
        # image is taller -> fit to height
        new_h = cell_h
        new_w = int(cell_h * img_aspect)

    # center within cell
    new_x = x + (cell_w - new_w) // 2
    new_y = y + (cell_h - new_h) // 2

    slide.shapes.add_picture(path, new_x, new_y, width=new_w, height=new_h)


def add_four_image_slide(prs, title_text, image_paths):
    blank_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.title.text = title_text
    tf = slide.shapes.title.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(20)

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    title_h = Inches(1.0)

    avail_w = slide_w
    avail_h = slide_h - title_h
    cell_w = avail_w // 2
    cell_h = avail_h // 2

    positions = {
        (1, "even"): (0,       title_h),
        (2, "even"): (cell_w,  title_h),
        (1, "odd"):  (0,       title_h + cell_h),
        (2, "odd"):  (cell_w,  title_h + cell_h),
    }

    for key, (x, y) in positions.items():
        path = image_paths.get(key)
        if path and os.path.exists(path):
            add_picture_fit(slide, path, x, y, cell_w, cell_h)
        else:
            tb = slide.shapes.add_textbox(x, y, cell_w, cell_h)
            tb.text_frame.text = f"MISSING:\n{path}"
    return slide


def resolve_path(template, endcap, charge, parity):
    fname = template.format(endcap=endcap, charge=charge, parity=parity)
    return os.path.join(plotdir, fname)


# ------------------------------------------------------------------
# Build the presentation
# ------------------------------------------------------------------
def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "GEM-CSC Alignment Plots",
        "Endcap 1/2 × Even/Odd chambers, split by muon charge",
    )

    for charge in charges:
        add_title_slide(prs, f"Charge: {charge}", "")

        for title, template, parity_split in plot_families:
            image_paths = {}
            for endcap in endcaps:
                for parity in parities:
                    if parity_split:
                        p = resolve_path(template, endcap, charge, parity)
                    else:
                        p = resolve_path(template, endcap, charge, parity="")
                    image_paths[(endcap, parity)] = p

            slide_title = f"data without alignment {title}  ({charge} endcap)"
            add_four_image_slide(prs, slide_title, image_paths)

    prs.save(output_pptx)
    print(f"Wrote {output_pptx}")


if __name__ == "__main__":
    build_pptx()